from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

def draw_symbol(list_error_slides, psr):
    if (list_error_slides):
        left = top = Inches(1.5)
        width = height = Inches(6)
        for i in list_error_slides:
            slide_error = psr.slides.get(i)
            chart_shape = slide_error.shapes.add_shape(MSO_SHAPE.MATH_MULTIPLY, left, top, width, height)
        
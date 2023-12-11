from Func_check_label.check_label import check_label
from Func_check_label.fill_color_textBox import fill_color_textBox
from Func_check_label.get_label_in_table import get_label_in_table
from Func_check_label.list_error import list_error
from Func_check_value.Render_textbox import renderData
from Func_check_value.getAll_textbox import get_all_textBox
from Func_check_value.getMenuTable import getTable
from Func_check_label.render_coor import render_arrow
from Func_check_label.render_coor import render_coor_label
from Func_check_label.classify_shapes import classify_kind_of_shape
from Func_check_label.find_label_connect import find_label_connect

from tkinter import messagebox
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from pptx import Presentation


# link_pp = 'C://Users//ADMIN//Desktop//code//python//Tkinter//formcheckpp//Dummy data//Dummy_pp.pptx'

def check_2(psr):
    # psr = Presentation(link)
    listSlides = psr.slides
        # get all text_box from slides
    data_textBox = get_all_textBox(listSlides, {})
        #_-----------------------------------
    (textbox_render, slide_my_table_id)  = renderData(data_textBox)
            #_-----------------------------------
    print('id slide menu: ',slide_my_table_id)
    slide = listSlides.get(slide_my_table_id)
    
    # find arrow shape and textbox
    (list_arrow,list_label) = classify_kind_of_shape(slide.shapes)
    list_arrow_render = render_arrow(list_arrow)
    list_label_render = render_coor_label(list_label)
    
    # render label cho de phan tich
    (label_arrows, render_label) = find_label_connect(list_arrow_render, list_label_render,list_label)
    print(render_label) #-----------------------data label render cuoi cung
    
    #select table       
    Menu_table = getTable(listSlides, slide_my_table_id)
    list_text_table = get_label_in_table(Menu_table)
    print("list_text_table: ", list_text_table) #-----------------------data text table cuoi cung
    
    # check label with table 
    if(len(list_text_table) > len(render_label)):
        messagebox.showinfo("Thông báo", "Số lượng label chi tiết không khớp với số chi tiết trong bảng tổng hợp")

        left = top = Inches(0.5)
        (width,height)= (Inches(5), Inches(0.5))
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textframe = textbox.text_frame
        textframe.text = "Số lượng label chỉ chi tiết đang ít hơn với số chi tiết trong bảng!"
        textframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        textframe.vertical_anchor = MSO_ANCHOR.MIDDLE
        textframe.word_wrap = True
        paragraph = textframe.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 0, 0)
        paragraph.font.size = Pt(25)
        paragraph.font.bold = True
        
    
    true_label = check_label(list_text_table,render_label, list_label)       
    # draw fill label error
    err_labels = list_error (label_arrows, true_label)  
    fill_color_textBox(err_labels)
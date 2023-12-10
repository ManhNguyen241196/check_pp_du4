from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

link_pp = 'C://Users//ADMIN//Desktop//code//python//Tkinter//formcheckpp//Dummy data//test.pptx'
psr = Presentation(link_pp)
listSlides = psr.slides

cur_slide = listSlides[0]
for shape in cur_slide.shapes:
    if shape.has_text_frame:
        textBox_value =  shape.text_frame
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 0, 0)  # red
psr.save("Dummy data/test_result.pptx")